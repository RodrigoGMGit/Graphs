#!/usr/bin/env python3
# generate_presentation.py
#
# Variante final:  Madurez & Dedicación centradas (70 % del ancho);
#                  TMD apilado, cada gráfico a 90 % del ancho (márgenes 0 .5″),
#                  bloque centrado verticalmente.
#                  Calidad sin cambios (rejilla 2×2).
# ---------------------------------------------------------------------------

import datetime as dt
import io
import os
import sys
from copy import deepcopy
from pathlib import Path
from typing import List, Optional, Tuple, cast

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches

from chapter_sync import graphs
from chapter_sync.file_processor import extract_date_from_standardized_filename

# ───── rutas
# Cuando se ejecuta desde un ejecutable PyInstaller, los recursos se
# descomprimen en ``sys._MEIPASS``. De lo contrario tomamos el directorio
# actual del script.
APP_DIR = Path(getattr(sys, "_MEIPASS", Path(__file__).resolve().parent))

# Plantilla de PowerPoint incluida en ``inputs/``
TEMPLATE_PATH = str(APP_DIR / "inputs" / "Template.pptx")

# Carpeta de salida junto al ejecutable (o al script durante el desarrollo)
if getattr(sys, "frozen", False):
    OUT_DIR = Path(sys.executable).resolve().parent / "outputs"
else:
    OUT_DIR = APP_DIR / "outputs"
os.makedirs(OUT_DIR, exist_ok=True)


# ───── capturar figuras
def capture(fn) -> List[io.BytesIO]:
    bufs: List[io.BytesIO] = []
    orig = plt.show

    def _cap(*a, **k):
        b = io.BytesIO()
        plt.savefig(b, format="png", dpi=150, bbox_inches="tight")
        b.seek(0)
        bufs.append(b)
        plt.close()

    plt.show = _cap
    fn()
    plt.show = orig
    return bufs


def _imgs(k, f) -> Tuple[List[io.BytesIO], Optional[str], Optional[dt.datetime]]:
    """Captura gráficos y devuelve buffers, ruta del archivo y fecha de corte.
    
    Returns:
        Tuple de (lista de buffers, ruta del archivo, fecha de corte)
    """
    p = graphs._resolve_path(None, k)
    if not p:
        return ([], None, None)
    
    buffers = capture(lambda: f(p))
    
    # Extraer fecha del nombre del archivo
    filename = os.path.basename(p)
    date_obj = extract_date_from_standardized_filename(filename)
    
    return (buffers, p, date_obj)


def main() -> None:
    # Check and download files if needed
    try:
        from chapter_sync.file_processor import check_and_download_if_needed

        check_and_download_if_needed(Path(graphs.FILES_DIR))
    except Exception as e:
        print(
            f"Error al verificar/descargar archivos: {e}. Continuando con archivos existentes."
        )

    imgs_mad, path_mad, date_mad = _imgs("madurez", graphs.plot_niveles_madurez)
    imgs_ded, path_ded, date_ded = _imgs("dedicacion", graphs.plot_dedicacion_tm)
    imgs_tmd, path_tmd, date_tmd = _imgs("tiempo", graphs.plot_tiempo_desarrollo)  # 2
    imgs_cal, path_cal, date_cal = _imgs("calidad", graphs.plot_calidad_pases)  # N

    prs = Presentation(TEMPLATE_PATH)
    SW: Emu = cast(Emu, prs.slide_width)
    SH: Emu = cast(Emu, prs.slide_height)

    PIC_W_STD = cast(Emu, int(SW * 0.70))  # 70 % ancho (Madurez/DR)
    PIC_W_TMD = cast(Emu, SW - Inches(1.0))  # 90 % ancho (0 .5″ márgenes)  # noqa: F841
    LEFT_STD = cast(Emu, (SW - PIC_W_STD) // 2)  # noqa: F841
    LEFT_TMD = Inches(0.5)  # noqa: F841
    TOP_MIN = Inches(0.8)
    GAP_V_TMD = Inches(0.40)  # noqa: F841

    def add_date_label(slide: Slide, date_obj: Optional[dt.datetime], file_path: Optional[str]) -> None:
        """Añade un cuadro de texto con la fecha de corte en la esquina superior derecha del slide."""
        if not date_obj and not file_path:
            return
        
        # Formatear la fecha
        if date_obj:
            date_str = date_obj.strftime("%d/%m/%Y")
            label_text = f"Fecha de corte: {date_str}"
        else:
            # Si no hay fecha, usar el nombre del archivo
            filename = os.path.basename(file_path) if file_path else "N/A"
            label_text = f"Archivo: {filename}"
        
        # Posición en la esquina superior derecha con márgenes
        width = Inches(3.0)  # Ancho suficiente para el texto
        left = cast(Emu, SW - width - Inches(0.5))  # Posicionado a la derecha
        top = Inches(0.3)
        height = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = label_text
        text_frame.word_wrap = False
        
        # Formato del texto
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Inches(0.12)  # ~9pt
        paragraph.font.name = "Calibri"
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Color negro
        paragraph.alignment = PP_ALIGN.RIGHT  # Right alignment

    def add_center(slide: Slide, buf: io.BytesIO, width: Emu) -> None:
        pic = slide.shapes.add_picture(buf, 0, 0, width)  # type: ignore
        pic.left = cast(Emu, (SW - pic.width) // 2)
        pic.top = cast(Emu, max(TOP_MIN, (SH - pic.height) // 2))

    if imgs_mad:
        add_center(prs.slides[2], imgs_mad[0], PIC_W_STD)
        add_date_label(prs.slides[2], date_mad, path_mad)

    # ——— Dedicación + Duración subtareas (mismo layout que TMD) ————
    if len(imgs_ded) >= 2:
        s3 = prs.slides[3]
        margin_h = Inches(0.5)
        gap_v = Inches(0.05)

        # => exactamente el mismo cálculo que TMD
        pic_w = cast(Emu, (SW - 2 * margin_h - Inches(0.25)) // 2)
        left_c = cast(Emu, (SW - pic_w) // 2)
        top_1 = Inches(0.2)

        shape1 = s3.shapes.add_picture(imgs_ded[0], left_c, top_1, pic_w)

        top_2 = cast(Emu, shape1.top + shape1.height + gap_v)
        s3.shapes.add_picture(imgs_ded[1], left_c, top_2, pic_w)
        add_date_label(s3, date_ded, path_ded)

    # Si por alguna razón solo llega una imagen (back-compat)
    elif imgs_ded:
        add_center(prs.slides[3], imgs_ded[0], PIC_W_STD)
        add_date_label(prs.slides[3], date_ded, path_ded)

    if len(imgs_tmd) >= 2:
        s5 = prs.slides[4]
        margin_h = Inches(0.5)
        gap_v = Inches(0.25)
        pic_w = cast(Emu, (SW - 2 * margin_h - Inches(0.25)) // 2)
        left_c = cast(Emu, (SW - pic_w) // 2)
        top_1 = Inches(1.0)
        shape1 = s5.shapes.add_picture(imgs_tmd[0], left_c, top_1, pic_w)
        top_2 = cast(Emu, shape1.top + shape1.height + gap_v)
        s5.shapes.add_picture(imgs_tmd[1], left_c, top_2, pic_w)
        add_date_label(s5, date_tmd, path_tmd)

    if imgs_cal:
        base = prs.slides[5]
        rect = (
            Inches(0.5),
            Inches(1.3),
            cast(Emu, SW - Inches(1.0)),
            cast(Emu, SH - Inches(1.8)),
        )
        idx, cur = 0, base
        while idx < len(imgs_cal):
            l, t, w, h = rect  # noqa: E741
            gap = Inches(0.15)
            cw, ch = cast(Emu, (w - gap) // 2), cast(Emu, (h - gap) // 2)
            for r in range(2):
                for c in range(2):
                    if idx >= len(imgs_cal):
                        break
                    cx = cast(Emu, l + c * (cw + gap))
                    cy = cast(Emu, t + r * (ch + gap))
                    cur.shapes.add_picture(imgs_cal[idx], cx, cy, cw, ch)
                    idx += 1
            # Añadir fecha de corte en cada slide de calidad
            add_date_label(cur, date_cal, path_cal)
            if idx < len(imgs_cal):
                new = prs.slides.add_slide(base.slide_layout)
                for shp in base.shapes:
                    if shp.is_placeholder and shp.text_frame:  # type: ignore
                        new.shapes._spTree.insert_element_before(
                            deepcopy(shp.element), "p:extLst"
                        )
                        break
                cur = new

    fname = dt.datetime.today().strftime("%Y-%m-%d_Presentation.pptx")
    prs.save(os.path.join(OUT_DIR, fname))
    print(f"\n✅ Presentación generada en outputs/{fname}\n")


if __name__ == "__main__":
    main()
