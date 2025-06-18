import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

from reportgen.config import Settings
from reportgen.graphs import bar_dedicacion, bar_tmd

logger = logging.getLogger(__name__)


# ── Context manager para capturar plots ───────────────────────────────────────
class capture_fig:
    """Captura la próxima figura creada por Matplotlib y la devuelve como BytesIO."""

    def __enter__(self):
        self._old_show = plt.show

        def _capture(*args, **kwargs):  # noqa: D401
            buf = BytesIO()
            plt.savefig(buf, format="png", dpi=200, bbox_inches="tight")
            buf.seek(0)
            self.buf = buf

        plt.show = _capture
        return self

    def __exit__(self, exc_type, exc, tb):
        plt.show = self._old_show


# ── Función principal ─────────────────────────────────────────────────────────
def build_presentation(settings: Settings, outfile: str | None = None) -> Path:
    prs = Presentation(settings.template_path)

    # ── Generar gráficos (usando las funciones del módulo graphs) ─────────────
    tmp_dir = Path("_tmp_figs")
    tg1 = bar_dedicacion(settings, tmp_dir)
    tg2 = bar_tmd(settings, tmp_dir)

    # ── Insertar en una slide vacía (ejemplo) ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco

    for idx, img_path in enumerate([tg1, tg2]):
        left = Inches(0.5)
        top = Inches(0.5 + idx * 3.5)
        height = Inches(3.0)
        slide.shapes.add_picture(str(img_path), left, top, height=height)

    # ── Guardar presentación ─────────────────────────────────────────────────
    out_dir = settings.output_dir
    out_dir.mkdir(parents=True, exist_ok=True)
    if not outfile:
        outfile = f"{datetime.now():%Y-%m-%d}_Presentation.pptx"
    out_path = out_dir / outfile

    prs.save(out_path)
    logger.info("Presentación guardada: %s", out_path)

    # limpiar temporales
    for p in tmp_dir.glob("*.png"):
        p.unlink()
    tmp_dir.rmdir()

    return out_path
