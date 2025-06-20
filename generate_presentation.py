#!/usr/bin/env python3
# generate_presentation.py
#
# Variante final:  Madurez & Dedicación centradas (70 % del ancho);
#                  TMD apilado, cada gráfico a 90 % del ancho (márgenes 0 .5″),
#                  bloque centrado verticalmente.
#                  Calidad sin cambios (rejilla 2×2).
# ---------------------------------------------------------------------------

import datetime as dt
import io
import os
from copy import deepcopy
from typing import List, cast

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Emu, Inches

import graphs
from utils import make_dirs_if_missing

assert graphs.PATHS_SET, "graphs.set_paths() must be called first"

# ───── rutas
TEMPLATE_PATH = r".\inputs\Template.pptx"
OUT_DIR = r".\outputs"
make_dirs_if_missing(OUT_DIR)


# ───── capturar figuras
def capture(fn) -> List[io.BytesIO]:
    bufs: List[io.BytesIO] = []
    orig = plt.show

    def _cap(*a, **k):
        b = io.BytesIO()
        plt.savefig(b, format="png", dpi=150, bbox_inches="tight")
        b.seek(0)
        bufs.append(b)
        plt.close()

    plt.show = _cap
    fn()
    plt.show = orig
    return bufs


def _imgs(k, f):
    p = graphs._resolve_path(None, k)
    return capture(lambda: f(p)) if p else []


imgs_mad = _imgs("madurez", graphs.plot_niveles_madurez)
imgs_ded = _imgs("dedicacion", graphs.plot_dedicacion_tm)
imgs_tmd = _imgs("tiempo", graphs.plot_tiempo_desarrollo)  # 2
imgs_cal = _imgs("calidad", graphs.plot_calidad_pases)  # N

# ───── plantilla
prs = Presentation(TEMPLATE_PATH)
SW: Emu = cast(Emu, prs.slide_width)
SH: Emu = cast(Emu, prs.slide_height)

# tamaños
PIC_W_STD = cast(Emu, int(SW * 0.70))  # 70 % ancho (Madurez/DR)
PIC_W_TMD = cast(Emu, SW - Inches(1.0))  # 90 % ancho (0 .5″ márgenes)
LEFT_STD = cast(Emu, (SW - PIC_W_STD) // 2)
LEFT_TMD = Inches(0.5)
TOP_MIN = Inches(0.8)
GAP_V_TMD = Inches(0.40)


# helper: centra gráfica con ancho fijo
def add_center(slide: Slide, buf: io.BytesIO, width: Emu) -> None:
    pic = slide.shapes.add_picture(buf, 0, 0, width)  # type: ignore[arg-type]
    pic.left = cast(Emu, (SW - pic.width) // 2)
    pic.top = cast(Emu, max(TOP_MIN, (SH - pic.height) // 2))


# ───── slide 3 Madurez
if imgs_mad:
    add_center(prs.slides[2], imgs_mad[0], PIC_W_STD)

# ───── slide 4 Dedicación
if imgs_ded:
    add_center(prs.slides[3], imgs_ded[0], PIC_W_STD)


# ───────── TMD – apilado sin estirar ─────────
if len(imgs_tmd) >= 2:
    s5 = prs.slides[4]

    # anchura igual a la usada antes en el formato lado-a-lado
    margin_h = Inches(0.5)
    gap_v = Inches(0.25)
    pic_w = cast(Emu, (SW - 2 * margin_h - Inches(0.25)) // 2)  # misma mitad de slide
    left_c = cast(Emu, (SW - pic_w) // 2)  # centrado
    top_1 = Inches(1.0)

    # primer gráfico
    shape1 = s5.shapes.add_picture(imgs_tmd[0], left_c, top_1, pic_w)  # type: ignore[arg-type]
    # segundo gráfico debajo, respetando gap_v
    top_2 = cast(Emu, shape1.top + shape1.height + gap_v)
    s5.shapes.add_picture(imgs_tmd[1], left_c, top_2, pic_w)  # type: ignore[arg-type]

# ───── calidad (igual que antes)
if imgs_cal:
    base = prs.slides[5]
    rect = (
        Inches(0.5),
        Inches(1.3),
        cast(Emu, SW - Inches(1.0)),
        cast(Emu, SH - Inches(1.8)),
    )
    idx, cur = 0, base
    while idx < len(imgs_cal):
        l, t, w, h = rect
        gap = Inches(0.15)
        cw, ch = cast(Emu, (w - gap) // 2), cast(Emu, (h - gap) // 2)
        for r in range(2):
            for c in range(2):
                if idx >= len(imgs_cal):
                    break
                cx = cast(Emu, l + c * (cw + gap))
                cy = cast(Emu, t + r * (ch + gap))
                cur.shapes.add_picture(imgs_cal[idx], cx, cy, cw, ch)  # type: ignore[arg-type]
                idx += 1
        if idx < len(imgs_cal):
            new = prs.slides.add_slide(base.slide_layout)
            for shp in base.shapes:
                if shp.is_placeholder and shp.text_frame:
                    new.shapes._spTree.insert_element_before(
                        deepcopy(shp.element), "p:extLst"
                    )
                    break
            cur = new

# ───── guardar
fname = dt.datetime.today().strftime("%Y-%m-%d_Presentation.pptx")
prs.save(os.path.join(OUT_DIR, fname))
print(f"\n✅ Presentación generada en outputs/{fname}\n")
